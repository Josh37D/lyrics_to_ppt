import re
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def parse_lyrics(lyrics):
    """
    Parse the lyrics into a structured format of slides and CODA.
    Each verse, chorus, and CODA is identified and stored as a whole paragraph.
    """
    paragraphs = lyrics.strip().split(
        "\n\n")  # Split lyrics into paragraphs by blank lines
    slides = []
    coda = None

    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        # Match verses like '1.', '2.', etc.
        if re.match(r'^\d+\.', paragraph):
            slides.append(("Verse", paragraph))
        # Match choruses like 'R:', 'R1:', 'R2:', etc.
        elif re.match(r'^R\d*:', paragraph):
            slides.append(("Chorus", paragraph))
        elif paragraph.startswith("C."):  # Match CODA
            coda = paragraph

    return slides, coda


def create_ppt(slides, coda, output_path):
    """
    Create a PowerPoint presentation from the parsed lyrics.
    Each verse, chorus, and CODA is placed on its own slide.
    """
    prs = Presentation()

    # Add slides for verses and choruses
    for slide_type, content in slides:
        slide = prs.slides.add_slide(
            prs.slide_layouts[6])  # Blank slide layout

        # Set slide background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

        # Remove all existing shapes from the slide
        for shape in list(slide.shapes):
            if hasattr(shape, "element"):
                shape.element.getparent().remove(shape.element)

        # Create a textbox that covers the entire slide
        textbox = slide.shapes.add_textbox(
            left=Inches(0), top=Inches(0), width=Inches(10), height=Inches(7.5)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)

        # Add only the lyrics content (no headers)
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(36)  # Large font size
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.alignment = PP_ALIGN.CENTER  # Center alignment

    # Add CODA slide if it exists
    if coda:
        slide = prs.slides.add_slide(
            prs.slide_layouts[6])  # Blank slide layout

        # Set slide background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

        # Remove all existing shapes from the slide
        for shape in list(slide.shapes):
            if hasattr(shape, "element"):
                shape.element.getparent().remove(shape.element)

        # Create a textbox that covers the entire slide
        textbox = slide.shapes.add_textbox(
            left=Inches(0), top=Inches(0), width=Inches(10), height=Inches(7.5)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)

        # Add only the CODA content (no header)
        p = text_frame.add_paragraph()
        p.text = coda
        p.font.size = Pt(36)  # Large font size
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.alignment = PP_ALIGN.CENTER  # Center alignment

    # Save the presentation
    prs.save(output_path)


def select_output_path():
    """Open a file dialog to select the output path."""
    file_path = filedialog.asksaveasfilename(
        defaultextension=".pptx",
        filetypes=[("PowerPoint Files", "*.pptx")],
        title="Save PowerPoint File"
    )
    output_path_var.set(file_path)


def generate_ppt():
    """Generate the PowerPoint presentation."""
    lyrics = lyrics_text.get("1.0", tk.END).strip()
    output_path = output_path_var.get()

    if not lyrics:
        messagebox.showerror(
            "Error", "Please paste the lyrics in the text box.")
        return

    if not output_path:
        messagebox.showerror("Error", "Please select an output path.")
        return

    slides, coda = parse_lyrics(lyrics)
    create_ppt(slides, coda, output_path)
    messagebox.showinfo(
        "Success", f"PowerPoint presentation created at:\n{output_path}")


# Create the main application window
root = tk.Tk()
root.title("Karaoke PPT Generator")

# Create and place the widgets
lyrics_label = tk.Label(root, text="Paste Lyrics:")
lyrics_label.pack(pady=5)

lyrics_text = tk.Text(root, wrap=tk.WORD, height=20, width=60)
lyrics_text.pack(pady=5)

output_path_label = tk.Label(root, text="Output Path:")
output_path_label.pack(pady=5)

output_path_var = tk.StringVar()
output_path_entry = tk.Entry(root, textvariable=output_path_var, width=50)
output_path_entry.pack(pady=5)

browse_button = tk.Button(root, text="Browse", command=select_output_path)
browse_button.pack(pady=5)

generate_button = tk.Button(root, text="Generate PPT", command=generate_ppt)
generate_button.pack(pady=20)

# Run the application
root.mainloop()
